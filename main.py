import os
import sys
import tkinter as tk
from tkinter import filedialog, messagebox
import pdfplumber
from pptx import Presentation
import openai
from pydub import AudioSegment
from pathlib import Path
import subprocess

# Redirect stdout and stderr to a log file
log_file = open(os.path.expanduser("~/SlideReader.log"), "w")
sys.stdout = log_file
sys.stderr = log_file

# Set up OpenAI API
openai.api_key = 'key'

# Define the path to the 'audios' directory in the user's home directory
audios_dir = os.path.expanduser("~/SlideReader_audios")

# Ensure 'audios' directory exists
os.makedirs(audios_dir, exist_ok=True)

def extract_text_from_pdf(file_path):
    with pdfplumber.open(file_path) as pdf:
        text = '\n'.join([page.extract_text() for page in pdf.pages if page.extract_text() is not None])
    return text

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

def chunk_text(text, max_tokens=4096):
    words = text.split()
    chunks = []
    current_chunk = []
    current_length = 0

    for word in words:
        word_length = len(word) + 1  # Adding 1 for the space
        if current_length + word_length > max_tokens:
            chunks.append(' '.join(current_chunk))
            current_chunk = []
            current_length = 0
        current_chunk.append(word)
        current_length += word_length

    if current_chunk:
        chunks.append(' '.join(current_chunk))

    return chunks

def text_to_speech_chunked(text, filename):
    chunks = chunk_text(text)
    audio_segments = []

    for i, chunk in enumerate(chunks):
        print(chunk)
        response = openai.audio.speech.create(
            model="tts-1-hd",
            voice="alloy",
            input=chunk,
        )
        audio_file_path = Path(audios_dir) / f"{filename}_{i}.mp3"
        response.write_to_file(audio_file_path)
        audio_segments.append(AudioSegment.from_mp3(audio_file_path))

    combined = sum(audio_segments)
    final_audio_path = Path(audios_dir) / f"{filename}.mp3"
    combined.export(final_audio_path, format="mp3")

    # Clean up temporary chunk files
    for i in range(len(chunks)):
        os.remove(Path(audios_dir) / f"{filename}_{i}.mp3")

    return final_audio_path

def handle_file_upload():
    file_path = filedialog.askopenfilename(
        filetypes=[("PDF files", "*.pdf"), ("PPTX files", "*.pptx")]
    )
    if not file_path:
        return

    file_extension = os.path.splitext(file_path)[1].lower()
    if file_extension == ".pdf":
        text = extract_text_from_pdf(file_path)
    elif file_extension == ".pptx":
        text = extract_text_from_pptx(file_path)
    else:
        messagebox.showerror("Error", "Unsupported file type")
        return

    if not text:
        messagebox.showerror("Error", "No text found in file")
        return

    filename = os.path.splitext(os.path.basename(file_path))[0]
    audio_file_path = text_to_speech_chunked(text, filename)
    messagebox.showinfo("Success", f"Audio file created: {audio_file_path}")
    refresh_audio_list()

def refresh_audio_list():
    listbox.delete(0, tk.END)
    for filename in os.listdir(audios_dir):
        listbox.insert(tk.END, filename)

def open_audio():
    selected_file = listbox.get(tk.ACTIVE)
    if selected_file:
        audio_file_path = os.path.join(audios_dir, selected_file)
        if os.name == 'nt':  # For Windows
            os.startfile(audio_file_path)
        elif os.name == 'posix':  # For macOS or Linux
            subprocess.call(('open', audio_file_path))

def delete_audio():
    selected_file = listbox.get(tk.ACTIVE)
    if selected_file:
        audio_file_path = os.path.join(audios_dir, selected_file)
        os.remove(audio_file_path)
        refresh_audio_list()

# Set up the Tkinter GUI
root = tk.Tk()
root.title("PDF/PPTX to Audio Converter")
root.geometry("500x400")

upload_button = tk.Button(root, text="Upload PDF/PPTX", command=handle_file_upload)
upload_button.pack(pady=10)

listbox = tk.Listbox(root)
listbox.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)

refresh_button = tk.Button(root, text="Refresh List", command=refresh_audio_list)
refresh_button.pack(pady=5)

open_button = tk.Button(root, text="Open Audio", command=open_audio)
open_button.pack(pady=5)

delete_button = tk.Button(root, text="Delete Audio", command=delete_audio)
delete_button.pack(pady=5)

refresh_audio_list()

# Ensure the application closes properly
def on_closing():
    root.destroy()

root.protocol("WM_DELETE_WINDOW", on_closing)
root.mainloop()

# Close the log file when done
log_file.close()
